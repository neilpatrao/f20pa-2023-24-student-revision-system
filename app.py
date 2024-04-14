import os
import pickle

import streamlit as st

from PyPDF2 import PdfReader
from pptx import Presentation
from dotenv import load_dotenv

from langchain_openai import OpenAI
from langchain_community.vectorstores.faiss import FAISS
from langchain_community.embeddings import HuggingFaceEmbeddings
from langchain_text_splitters import CharacterTextSplitter
from langchain.chains.question_answering import load_qa_chain

class Types:
    PDF = "application/pdf"
    PPTX = "application/vnd.openxmlformats-officedocument.presentationml.presentation"
    DOCX = "application/vnd.openxmlformats-officedocument.wordprocessingml.document"

def read_pdf(file):
    text = ""
    pdf = PdfReader(file)
    for page in pdf.pages:
        text += page.extract_text()
    return text

def read_pptx(file):
    text = ""
    prs = Presentation(file)
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text
    return text


def get_raw_text(docs):
    text = ""
    for doc in docs:
        match doc.type:
            case Types.PDF:
                print("\nPDF\n")
                text += read_pdf(doc)
            case Types.PPTX:
                print("\nPPTX\n")
                text += read_pptx(doc)
    return text

def get_text_chunks(text):
    text_splitter = CharacterTextSplitter(
        separator="\n",
        chunk_size=1000,
        chunk_overlap=200,
        length_function=len
	)
    chunks = text_splitter.split_text(text)
    return chunks

def create_vectorstore(text_chunks):
    embeddings = HuggingFaceEmbeddings()
    vectorstore = FAISS.from_texts(texts=text_chunks, embedding=embeddings)
    return vectorstore

def create_chain():
    llm = OpenAI(base_url="http://localhost:1234/v1", api_key="lm-studio")
    # llm = OpenAI()
    chain = load_qa_chain(llm=llm, chain_type="stuff")
    return chain

def main():
    load_dotenv()
    st.set_page_config(page_title="Chat with multiple PDFs", page_icon=":books:")

    if "chain" not in st.session_state:
        st.session_state.chain = None

    if "vectorstore" not in st.session_state:
        st.session_state.vectorstore = None

    st.header("Student Revision System :books:")	

    with st.sidebar:
        st.subheader("Your Documents")
        pdfs = st.file_uploader("Upload your PDFs", type=["pdf", "pptx", "docx"], accept_multiple_files=True)
        print(pdfs)
        if st.button("Process"):
            with st.spinner("Processing Documents"):
                # Get raw text from documents
                raw_text = get_raw_text(pdfs)
                text_chunks = get_text_chunks(raw_text)
                st.session_state.vectorstore = create_vectorstore(text_chunks)
                st.session_state.chain = create_chain()
    
    # query = st.text_input("Ask your question")
    query = st.chat_input("Ask your question")
    if query:
        with st.chat_message("user"):
            st.write(query)
        docs = st.session_state.vectorstore.similarity_search(query=query)
        print(f"\n{docs}\n")
        response = st.session_state.chain.run(input_documents=docs, question=query)
        with st.chat_message("assistant"):
            st.write(response)


if __name__ == "__main__":
    main()